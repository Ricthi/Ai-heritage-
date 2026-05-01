import streamlit as st
from pptx import Presentation

def extract_pptx_info_from_file(file):
    """Read a PPTX file-like object and show slide content in Streamlit."""
    prs = Presentation(file)
    st.write("--- Presentation Info ---")

    for i, slide in enumerate(prs.slides):
        st.write(f"\n### Slide {i + 1}")
        for shape in slide.shapes:
            # Text content
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    st.write(f"- {text}")
            # 13 = Picture in python-pptx
            if getattr(shape, "shape_type", None) == 13:
                st.write("- [Image Present]")

def main():
    st.title("Tamil PPTX Slide Inspector")

    uploaded = st.file_uploader("Upload a PPTX file", type=["pptx"])
    if uploaded is not None:
        extract_pptx_info_from_file(uploaded)
    else:
        st.info("Please upload a .pptx file to analyze.")

if __name__ == "__main__":
    main()
